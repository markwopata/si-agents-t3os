from __future__ import annotations

from io import BytesIO
from pathlib import Path

import requests
from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/mark.wopata/Documents/projects/SI-agent")
OUTPUT = ROOT / "docs" / "equipmentshare-analytics-agent-executive-deck-v1.pptx"

GRAPHITE = RGBColor(0x1F, 0x2A, 0x37)
SAND = RGBColor(0xF4, 0xEF, 0xE6)
GOLD = RGBColor(0xD8, 0xA7, 0x3B)
SLATE = RGBColor(0x5B, 0x64, 0x70)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEC, 0xE7, 0xDD)
MUTED = RGBColor(0x7B, 0x83, 0x8D)
SOFT_GOLD = RGBColor(0xF1, 0xDB, 0xA4)

TITLE_FONT = "Aptos Display"
BODY_FONT = "Aptos"

ARCH_DIAGRAM_URL = (
    "https://s3-alpha.figma.com/thumbnails/b9ccdb88-ea5b-4958-bdb2-25662da30f7a"
    "?X-Amz-Algorithm=AWS4-HMAC-SHA256&X-Amz-Credential=AKIAQ4GOSFWCYKED6IIG%2F"
    "20260406%2Fus-west-2%2Fs3%2Faws4_request&X-Amz-Date=20260406T124528Z&"
    "X-Amz-Expires=604800&X-Amz-SignedHeaders=host&X-Amz-Signature="
    "ccd73e57fad7c56c9149c40d5e8a0573dff4055d5884658b980c4071db0f8090"
)
LIFECYCLE_DIAGRAM_URL = (
    "https://s3-alpha.figma.com/thumbnails/7d2da1ad-9988-469c-9f4f-4cff683479f5"
    "?X-Amz-Algorithm=AWS4-HMAC-SHA256&X-Amz-Credential=AKIAQ4GOSFWCYKED6IIG%2F"
    "20260406%2Fus-west-2%2Fs3%2Faws4_request&X-Amz-Date=20260406T124539Z&"
    "X-Amz-Expires=604800&X-Amz-SignedHeaders=host&X-Amz-Signature="
    "b01fcd037b41b8776a8d1491f02fdb96b5cc74afcdc15a0e15686a5214101b6f"
)


def download_image(url: str) -> BytesIO | None:
    try:
        response = requests.get(url, timeout=30)
        response.raise_for_status()
    except Exception:
        return None
    try:
        img = BytesIO(response.content)
        Image.open(img).verify()
        img.seek(0)
        return img
    except Exception:
        return None


def set_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, *, font_size=20, color=GRAPHITE,
                bold=False, name=BODY_FONT, align=PP_ALIGN.LEFT, linespacing=1.1):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.space_after = Pt(2)
    p.line_spacing = linespacing
    run = p.runs[0]
    run.font.name = name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, *, font_size=18, color=GRAPHITE,
                bullet_color=GOLD, bg=None):
    if bg is not None:
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = bg
        card.line.color.rgb = bg
    box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), height - Inches(0.2))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.bullet = True
        p.space_after = Pt(8)
        p.line_spacing = 1.15
        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
            run.text = item
        run.font.name = BODY_FONT
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return box


def add_title(slide, title: str, subtitle: str | None = None, *, dark=False) -> None:
    color = WHITE if dark else GRAPHITE
    accent = SOFT_GOLD if dark else GOLD
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(0.45), Inches(1.15), Inches(0.1)
    ).fill.solid()
    band = slide.shapes[-1]
    band.fill.fore_color.rgb = accent
    band.line.color.rgb = accent
    add_textbox(
        slide, Inches(0.65), Inches(0.72), Inches(8.2), Inches(0.8), title,
        font_size=28, color=color, bold=True, name=TITLE_FONT
    )
    if subtitle:
        add_textbox(
            slide, Inches(0.65), Inches(1.28), Inches(10.5), Inches(0.6), subtitle,
            font_size=14, color=SOFT_GOLD if dark else SLATE, name=BODY_FONT
        )


def add_footer(slide, text: str, *, dark=False) -> None:
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(7.14), Inches(13.33), Inches(0.36)
    )
    band = slide.shapes[-1]
    band.fill.solid()
    band.fill.fore_color.rgb = GRAPHITE if not dark else RGBColor(0x14, 0x1C, 0x26)
    band.line.color.rgb = band.fill.fore_color.rgb
    add_textbox(
        slide, Inches(0.72), Inches(7.18), Inches(12), Inches(0.2), text,
        font_size=10, color=WHITE if not dark else LIGHT
    )


def add_metric_card(slide, left, top, width, height, value: str, label: str, subtitle: str):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT
    add_textbox(slide, left + Inches(0.18), top + Inches(0.14), width - Inches(0.36), Inches(0.38),
                value, font_size=22, color=GRAPHITE, bold=True, name=TITLE_FONT)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.56), width - Inches(0.36), Inches(0.25),
                label, font_size=11, color=GOLD, bold=True)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.82), width - Inches(0.36), height - Inches(0.96),
                subtitle, font_size=11, color=SLATE)


def add_domain_chip(slide, left, top, width, text):
    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.38)
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = LIGHT
    chip.line.color.rgb = LIGHT
    add_textbox(slide, left + Inches(0.12), top + Inches(0.08), width - Inches(0.24), Inches(0.18),
                text, font_size=11, color=GRAPHITE, bold=True)


def add_chart(slide):
    chart_data = CategoryChartData()
    chart_data.categories = [
        "Inventory Information",
        "Company Directory",
        "Customer Dashboard",
        "Salesperson",
        "Market Dashboard",
        "Service Dashboard",
    ]
    chart_data.add_series("Runs", (227990, 86447, 78887, 73961, 58570, 55349))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(6.55), Inches(1.75), Inches(6.0), Inches(4.2), chart_data
    ).chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.reverse_order = True
    series = chart.series[0]
    fill = series.format.fill
    fill.solid()
    fill.fore_color.rgb = GOLD
    series.format.line.color.rgb = GOLD


def build_cover(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, GRAPHITE)
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(5.7), Inches(13.33), Inches(1.8)
    )
    band = slide.shapes[-1]
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(0x16, 0x1F, 0x2B)
    band.line.color.rgb = band.fill.fore_color.rgb
    for left, top, w, h, color in [
        (0.65, 0.55, 2.35, 0.16, GOLD),
        (9.35, 0.9, 2.75, 1.25, RGBColor(0x2C, 0x3A, 0x49)),
        (10.2, 2.0, 1.6, 0.95, GOLD),
        (8.9, 2.95, 2.8, 1.3, RGBColor(0x39, 0x46, 0x54)),
    ]:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
    add_textbox(
        slide, Inches(0.78), Inches(1.1), Inches(7.0), Inches(1.4),
        "EquipmentShare Intelligence Platform", font_size=30, color=WHITE, bold=True, name=TITLE_FONT
    )
    add_textbox(
        slide, Inches(0.78), Inches(2.02), Inches(6.8), Inches(1.1),
        "Analytics Agent v1", font_size=24, color=SOFT_GOLD, bold=True, name=TITLE_FONT
    )
    add_textbox(
        slide, Inches(0.82), Inches(3.0), Inches(6.25), Inches(1.15),
        "An analytics-first company intelligence platform grounded in repo-aware reasoning, governed execution, and real company usage signals.",
        font_size=16, color=LIGHT
    )
    add_textbox(
        slide, Inches(0.82), Inches(5.98), Inches(8.0), Inches(0.38),
        "Executive working draft built from prototype results, Slack discovery, and LookML usage analysis.",
        font_size=12, color=LIGHT
    )


def build_why_now(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, SAND)
    add_title(slide, "Why This Needs To Exist", "The highest-value analytics work is investigative, cross-system, and governed.")
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.75), Inches(4.35), Inches(4.85)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = GRAPHITE
    panel.line.color.rgb = GRAPHITE
    add_textbox(
        slide, Inches(1.0), Inches(2.0), Inches(3.75), Inches(0.8),
        "Generic LLMs can summarize. They cannot reliably reason across EquipmentShare code, semantics, freshness, and permissions.",
        font_size=20, color=WHITE, bold=True, name=TITLE_FONT
    )
    add_textbox(
        slide, Inches(1.0), Inches(3.45), Inches(3.5), Inches(1.6),
        "This platform is meant to close that gap by combining trusted code context, governed execution, and domain-aware routing.",
        font_size=14, color=LIGHT
    )
    items = [
        "Important questions are often “why did this move?” or “why doesn’t this tie out?” rather than “what is metric X?”",
        "Knowledge is fragmented across dbt repos, LookML, docs, dashboards, operational replicas, and analyst memory.",
        "Freshness, provenance, source trust, and sensitivity handling are first-order requirements.",
        "The same backbone needs to serve analysts, operators, internal apps, and future APIs."
    ]
    add_bullets(slide, Inches(5.35), Inches(1.75), Inches(7.2), Inches(4.85), items, bg=WHITE)
    add_footer(slide, "EquipmentShare analytics-agent proposal")


def build_proof(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_title(slide, "Proof From Real Prototype Work", "The prototype already works across multiple business domains and artifact types.")
    cards = [
        ("Branch earnings", "Investigations, invoice logic, variance drivers"),
        ("Fleet and OEC", "Utilization, unavailable OEC, asset-level analysis"),
        ("Pricing", "Rate achievement, benchmark and floor rate analysis"),
        ("TCO", "Ownership economics, depreciation, Rouse comparisons"),
        ("Materials", "BiSTrack and Sage validation, location-month revenue"),
        ("People", "Base pay, bonus, fringe, payroll variance"),
        ("Training", "Docebo completion and historical progress"),
        ("Valuation", "NBV, wholesale value, OWN and payout logic"),
    ]
    lefts = [0.75, 3.95, 7.15, 10.35]
    tops = [1.95, 4.1]
    i = 0
    for top in tops:
        for left in lefts:
            title, body = cards[i]
            card = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(2.55), Inches(1.65)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = SAND
            card.line.color.rgb = LIGHT
            add_textbox(slide, Inches(left + 0.16), Inches(top + 0.16), Inches(2.15), Inches(0.3),
                        title, font_size=13, color=GRAPHITE, bold=True)
            accent = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left + 0.16), Inches(top + 0.52), Inches(0.55), Inches(0.07)
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = GOLD
            accent.line.color.rgb = GOLD
            add_textbox(slide, Inches(left + 0.16), Inches(top + 0.68), Inches(2.15), Inches(0.7),
                        body, font_size=10.5, color=SLATE)
            i += 1
    add_footer(slide, "Representative solved work across finance, fleet, pricing, materials, people, training, and valuation")


def build_usage(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, SAND)
    add_title(slide, "What Company Usage Is Telling Us", "Slack discovery and Looker usage point to the first domains the agent must understand.")
    add_metric_card(slide, Inches(0.75), Inches(1.75), Inches(2.7), Inches(1.35), "227,990", "Inventory Information", "Largest observed Looker usage signal in the export.")
    add_metric_card(slide, Inches(3.65), Inches(1.75), Inches(2.7), Inches(1.35), "86,447", "Company Directory", "Shared semantic and lookup behavior matters.")
    add_metric_card(slide, Inches(0.75), Inches(3.2), Inches(2.7), Inches(1.35), "78,887", "Customer Dashboard", "Commercial and account analytics are core.")
    add_metric_card(slide, Inches(3.65), Inches(3.2), Inches(2.7), Inches(1.35), "73,961", "Salesperson", "Revenue attribution and performance are heavily used.")
    add_metric_card(slide, Inches(0.75), Inches(4.65), Inches(2.7), Inches(1.35), "58,570", "Market Dashboard", "Branch and market performance remain central.")
    add_metric_card(slide, Inches(3.65), Inches(4.65), Inches(2.7), Inches(1.35), "55,349", "Service Dashboard", "Service and maintenance deserve their own domain.")
    add_chart(slide)
    add_textbox(
        slide, Inches(6.6), Inches(6.12), Inches(5.85), Inches(0.48),
        "Slack validated live language such as branch earnings, OEC, utilization, rate achievement, Anaplan, OWN, TCO, BiSTrack, and Intacct.",
        font_size=10.5, color=SLATE
    )
    add_footer(slide, "Usage evidence suggests fleet/assets, customer/revenue, branch performance, maintenance, and shared semantics are first-class priorities")


def build_domain(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_title(slide, "Domain Backbone Plus Shared Semantic Layer", "Business domains sit underneath a conformed layer used across the entire analytics surface.")
    layer = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.6), Inches(11.25), Inches(0.8)
    )
    layer.fill.solid()
    layer.fill.fore_color.rgb = GRAPHITE
    layer.line.color.rgb = GRAPHITE
    add_textbox(
        slide, Inches(1.3), Inches(1.88), Inches(10.7), Inches(0.25),
        "Shared semantic layer: market and region xwalks, customer hierarchies, fiscal calendars, asset hierarchies, org structures",
        font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER
    )
    domains = [
        "Branch earnings", "General ledger", "Fixed assets", "Customers and revenue",
        "Pricing", "Fleet and utilization", "Maintenance", "TCO",
        "Materials", "People and payroll", "Learning", "Planning",
        "OWN program", "Asset disposition",
    ]
    lefts = [0.95, 3.2, 5.45, 7.7, 9.95]
    tops = [3.0, 3.9, 4.8]
    idx = 0
    for top in tops:
        for left in lefts:
            if idx >= len(domains):
                break
            add_domain_chip(slide, Inches(left), Inches(top), Inches(2.0), domains[idx])
            idx += 1
    note = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(6.0), Inches(11.2), Inches(0.72)
    )
    note.fill.solid()
    note.fill.fore_color.rgb = SAND
    note.line.color.rgb = SAND
    add_textbox(
        slide, Inches(1.25), Inches(6.23), Inches(10.8), Inches(0.25),
        "OWN and asset disposition should be modeled separately. Shared lookup assets like Company Directory and market_region_xwalk show why the semantic layer cannot be buried inside one domain.",
        font_size=11.5, color=SLATE, align=PP_ALIGN.CENTER
    )
    add_footer(slide, "Domains should be business-shaped, with conformed dimensions above them")


def build_metrics(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, SAND)
    add_title(slide, "Most-Used Metric Families", "The first metric registry should begin with the families that recur in high-usage dashboards and matched LookML.")
    columns = [
        ("Fleet and assets", ["inventory counts", "asset status", "asset availability", "ttl_oec", "rental_revenue", "fin_util"]),
        ("Customer and commercial", ["customer revenue", "salesperson performance", "commission_total", "final_commission_payout", "clawback_total", "quote-request flow"]),
        ("Branch and accounting", ["branch_earnings_amount", "gaap_amount", "difference", "collections activity", "AP and card operations"]),
        ("Maintenance and pricing", ["sum_of_amount", "parts activity", "warranty activity", "total_points", "benchmark %", "count_of_assets"]),
    ]
    lefts = [0.7, 3.95, 7.2, 10.0]
    widths = [2.95, 2.95, 2.55, 2.55]
    for (title, items), left, width in zip(columns, lefts, widths):
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.85), Inches(width), Inches(4.95)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LIGHT
        add_textbox(slide, Inches(left + 0.18), Inches(2.05), Inches(width - 0.36), Inches(0.4),
                    title, font_size=14, color=GRAPHITE, bold=True)
        add_bullets(
            slide, Inches(left + 0.08), Inches(2.42), Inches(width - 0.16), Inches(4.1), items,
            font_size=11, color=SLATE
        )
    add_footer(slide, "Metric families sourced from the Looker dashboard usage export and direct LookML inspection")


def build_architecture(prs: Presentation, img_data: BytesIO | None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_title(slide, "Platform Architecture", "The platform routes messy business questions through trusted knowledge and governed execution before assembling a response.")
    bullets = [
        "Gateway and orchestration layer manages request flow.",
        "Domain router decides whether the question belongs to branch performance, accounting, fleet, maintenance, pricing, or another domain.",
        "Knowledge services rank code, docs, reviewed memory, and shared semantic infrastructure.",
        "Execution services handle read-only SQL, freshness checks, provenance, and permission policy.",
        "Responses can be explanations, SQL, direct answers, or artifacts."
    ]
    add_bullets(slide, Inches(0.75), Inches(1.8), Inches(4.4), Inches(4.9), bullets, bg=SAND, font_size=12)
    if img_data is not None:
        slide.shapes.add_picture(img_data, Inches(5.45), Inches(1.8), width=Inches(7.1), height=Inches(4.95))
    else:
        panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.45), Inches(1.8), Inches(7.1), Inches(4.95)
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = SAND
        panel.line.color.rgb = LIGHT
        add_textbox(
            slide, Inches(5.85), Inches(3.0), Inches(6.3), Inches(0.55),
            "Architecture diagram unavailable", font_size=22, color=GRAPHITE, bold=True, name=TITLE_FONT, align=PP_ALIGN.CENTER
        )
        add_textbox(
            slide, Inches(5.95), Inches(3.7), Inches(6.1), Inches(1.2),
            "The Figma-generated diagram URL was no longer returning a valid image at build time, so this slide keeps the architecture narrative and can be refreshed later with the visual asset.",
            font_size=12, color=SLATE, align=PP_ALIGN.CENTER
        )
    add_footer(slide, "Architecture visual adapted from the Figma-generated analytics-agent map")


def build_lifecycle(prs: Presentation, img_data: BytesIO | None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, GRAPHITE)
    add_title(slide, "Question-To-Answer Lifecycle", "Governance is built into the process, not bolted on afterward.", dark=True)
    if img_data is not None:
        slide.shapes.add_picture(img_data, Inches(0.75), Inches(1.7), width=Inches(7.1), height=Inches(4.9))
    else:
        panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.7), Inches(7.1), Inches(4.9)
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor(0x28, 0x35, 0x44)
        panel.line.color.rgb = RGBColor(0x28, 0x35, 0x44)
        add_textbox(
            slide, Inches(1.15), Inches(3.0), Inches(6.3), Inches(0.55),
            "Lifecycle diagram unavailable", font_size=22, color=WHITE, bold=True, name=TITLE_FONT, align=PP_ALIGN.CENTER
        )
        add_textbox(
            slide, Inches(1.1), Inches(3.7), Inches(6.4), Inches(1.2),
            "The live diagram image could not be embedded during build, so this version preserves the lifecycle explanation and can be refreshed later with the visual asset.",
            font_size=12, color=LIGHT, align=PP_ALIGN.CENTER
        )
    items = [
        "Interpret intent and detect domain",
        "Pull likely logic from dbt, LookML, docs, and reviewed memory",
        "Choose sources by trust, grain, freshness, and sensitivity",
        "Run governed read-only SQL only when live data is required",
        "Validate with provenance and known caveats",
        "Capture reusable prompt and query patterns back into reviewed memory",
    ]
    add_bullets(slide, Inches(8.1), Inches(1.72), Inches(4.5), Inches(4.92), items,
                font_size=11.5, color=WHITE, bg=RGBColor(0x28, 0x35, 0x44))
    add_footer(slide, "Lifecycle visual adapted from the Figma-generated question flow", dark=True)


def build_governance(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, SAND)
    add_title(slide, "Governance, Trust, And Sensitivity", "A durable company capability needs explicit trust tiers and permission behavior.")
    tiers = [
        ("Tier 1", "Canonical dbt marts and official docs", GRAPHITE, WHITE),
        ("Tier 2", "LookML and downstream semantic logic", SLATE, WHITE),
        ("Tier 3", "Curated playbooks and reviewed memory", GOLD, GRAPHITE),
        ("Tier 4", "Staging and raw operational content", LIGHT, GRAPHITE),
    ]
    top = 2.0
    widths = [5.6, 4.65, 3.65, 2.7]
    lefts = [0.95, 1.43, 1.93, 2.42]
    heights = [0.85, 0.75, 0.7, 0.65]
    for (name, desc, fill_color, text_color), left, width, height in zip(tiers, lefts, widths, heights):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.rotation = 180
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = fill_color
        add_textbox(slide, Inches(left + 0.3), Inches(top + 0.18), Inches(width - 0.6), Inches(0.3),
                    f"{name}: {desc}", font_size=12, color=text_color, bold=True, align=PP_ALIGN.CENTER)
        top += 0.68
    add_textbox(
        slide, Inches(7.15), Inches(2.0), Inches(5.1), Inches(0.5),
        "Sensitivity tiers", font_size=20, color=GRAPHITE, bold=True, name=TITLE_FONT
    )
    tiers_right = [
        "broad_internal",
        "operational_sensitive",
        "finance_restricted",
        "confidential_people",
        "customer_sensitive",
    ]
    y = 2.75
    for item in tiers_right:
        add_domain_chip(slide, Inches(7.2), Inches(y), Inches(4.5), item)
        y += 0.63
    add_textbox(
        slide, Inches(7.2), Inches(6.0), Inches(5.0), Inches(0.55),
        "Every answer should surface freshness, provenance, and the source trust logic that produced it.",
        font_size=12, color=SLATE
    )
    add_footer(slide, "Trust tiers and sensitivity behavior are part of the product, not back-office policy only")


def build_close(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_title(slide, "Prompt Portfolio And 90-Day Build Plan", "Launch with real prompts, real domains, and a pilot that can become the analytics-agent foundation.")
    left = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.78), Inches(6.0), Inches(4.95)
    )
    left.fill.solid()
    left.fill.fore_color.rgb = SAND
    left.line.color.rgb = SAND
    add_textbox(slide, Inches(1.0), Inches(2.0), Inches(5.3), Inches(0.35),
                "Prompt portfolio", font_size=18, color=GRAPHITE, bold=True, name=TITLE_FONT)
    prompts = [
        "Why did branch earnings move this month, and which categories drove the variance?",
        "Which customers are driving revenue growth, and which top accounts are the biggest movers by branch?",
        "Which markets are underperforming on utilization, and which classes are driving the highest maintenance cost or downtime?",
        "Explain how rate achievement was calculated for this invoice.",
        "How did this transaction hit the GL, and why does this reconciliation not tie out?",
        "Walk me through the official TCO formula and separate OWN economics from asset disposition performance.",
    ]
    add_bullets(slide, Inches(0.95), Inches(2.42), Inches(5.5), Inches(4.0), prompts, font_size=10.6)
    add_textbox(slide, Inches(7.15), Inches(2.0), Inches(4.5), Inches(0.35),
                "90-day build plan", font_size=18, color=GRAPHITE, bold=True, name=TITLE_FONT)
    phases = [
        ("Phase 1", "Stand up the analyst-facing ask surface, source registry, and domain routing."),
        ("Phase 2", "Tighten approved Slack ingestion, prompt persistence, and the shared semantic layer."),
        ("Phase 3", "Expand verified query coverage plus freshness, provenance, and sensitivity services."),
    ]
    y = 2.6
    for phase, desc in phases:
        circle = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, Inches(7.2), Inches(y), Inches(0.38), Inches(0.38)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = GOLD
        circle.line.color.rgb = GOLD
        add_textbox(slide, Inches(7.72), Inches(y - 0.02), Inches(4.3), Inches(0.22),
                    phase, font_size=12.5, color=GRAPHITE, bold=True)
        add_textbox(slide, Inches(7.72), Inches(y + 0.22), Inches(4.3), Inches(0.55),
                    desc, font_size=11.5, color=SLATE)
        y += 1.2
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.15), Inches(5.95), Inches(4.8), Inches(0.72)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = GRAPHITE
    band.line.color.rgb = GRAPHITE
    add_textbox(slide, Inches(7.45), Inches(6.18), Inches(4.2), Inches(0.25),
                "Decision: approve an analytics-first pilot and turn it into the dedicated analytics-agent project.",
                font_size=11.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "First working deck generated locally because the Figma Slides endpoint was unavailable")


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    arch = download_image(ARCH_DIAGRAM_URL)
    lifecycle = download_image(LIFECYCLE_DIAGRAM_URL)

    build_cover(prs)
    build_why_now(prs)
    build_proof(prs)
    build_usage(prs)
    build_domain(prs)
    build_metrics(prs)
    build_architecture(prs, arch)
    build_lifecycle(prs, lifecycle)
    build_governance(prs)
    build_close(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    main()
